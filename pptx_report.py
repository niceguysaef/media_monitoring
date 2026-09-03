from __future__ import annotations

from datetime import datetime
from io import BytesIO
import re
from typing import Any, Dict, Iterable, List, Optional, Sequence
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_WIDTH = 13.333333
SLIDE_HEIGHT = 7.5
PX = 1 / 96

C = {
    "dark": "24150E",
    "muted": "7C6256",
    "orange_dark": "B83A18",
    "orange": "ED6E1F",
    "gold": "ECA81A",
    "soft": "FFF1E6",
    "surface": "FFFAF6",
    "line": "ECD6C7",
    "positive": "287A55",
    "neutral": "6C78A2",
    "negative": "C84B38",
    "white": "FFFFFF",
}


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.lstrip("#"))


def _safe_number(value: Any) -> float:
    try:
        return float(value or 0)
    except (TypeError, ValueError):
        return 0.0


def _compact(value: Any) -> str:
    number = _safe_number(value)
    absolute = abs(number)
    if absolute >= 1_000_000_000:
        return f"{number / 1_000_000_000:.1f}B".replace(".0B", "B")
    if absolute >= 1_000_000:
        return f"{number / 1_000_000:.1f}M".replace(".0M", "M")
    if absolute >= 1_000:
        return f"{number / 1_000:.1f}K".replace(".0K", "K")
    return f"{number:,.0f}"


def _truncate(value: Any, maximum: int) -> str:
    text = " ".join(str(value or "").split())
    return text if len(text) <= maximum else f"{text[: maximum - 1].rstrip()}…"


def _short_date(value: str, malay: bool) -> str:
    try:
        parsed = datetime.strptime(value, "%Y-%m-%d")
        months = ["Jan", "Feb", "Mac", "Apr", "Mei", "Jun", "Jul", "Ogo", "Sep", "Okt", "Nov", "Dis"]
        return f"{parsed.day} {months[parsed.month - 1]}" if malay else parsed.strftime("%-d %b")
    except (TypeError, ValueError):
        return str(value or "")


def _source_label(value: Any, malay: bool) -> str:
    normalized = str(value or "").lower()
    if normalized == "twitter":
        return "X / Twitter"
    if normalized == "youtube":
        return "YouTube"
    if normalized == "tiktok":
        return "TikTok"
    if normalized:
        return normalized[0].upper() + normalized[1:]
    return "Lain-lain" if malay else "Other"


def _set_background(slide: Any, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _shape(
    slide: Any,
    kind: MSO_SHAPE,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: str,
    line: Optional[str] = None,
) -> Any:
    item = slide.shapes.add_shape(
        kind,
        Inches(left * PX),
        Inches(top * PX),
        Inches(width * PX),
        Inches(max(height, 0.5) * PX),
    )
    item.fill.solid()
    item.fill.fore_color.rgb = _rgb(fill)
    item.line.color.rgb = _rgb(line or fill)
    item.line.width = Pt(0.75 if line else 0)
    return item


def _rule(slide: Any, left: float, top: float, width: float, color: str) -> None:
    item = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left * PX),
        Inches(top * PX),
        Inches(width * PX),
        Pt(0.75),
    )
    item.fill.solid()
    item.fill.fore_color.rgb = _rgb(color)
    item.line.fill.background()


def _text(
    slide: Any,
    value: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: float = 18,
    bold: bool = False,
    color: str = C["dark"],
    align: str = "left",
    valign: str = "top",
    margin: float = 0,
) -> Any:
    box = slide.shapes.add_textbox(
        Inches(left * PX), Inches(top * PX), Inches(width * PX), Inches(height * PX)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = Inches(margin * PX)
    frame.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }[valign]
    lines = str(value if value is not None else "").split("\n")
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }[align]
        paragraph.space_before = paragraph.space_after = Pt(0)
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size * 0.75)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = _rgb(color)
    return box


def _header(slide: Any, title: str, kicker: str, page: int) -> None:
    _set_background(slide, C["surface"])
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 1280, 116, C["orange_dark"])
    _text(slide, kicker.upper(), 64, 22, 820, 22, size=16, bold=True, color="FFD9BD")
    _text(slide, title, 64, 49, 1080, 50, size=36, bold=True, color=C["white"])
    _text(slide, str(page).zfill(2), 1170, 55, 54, 28, size=18, bold=True, color=C["white"], align="right")


def _footer(slide: Any, dataset: Dict[str, Any], page: int, t: Any) -> None:
    _rule(slide, 64, 682, 1152, C["line"])
    period = f"{dataset['dateRange']['from']} – {dataset['dateRange']['to']}"
    _text(slide, f"{dataset['project']['name']}  •  {period}", 64, 690, 850, 20, size=12, color=C["muted"])
    _text(slide, t("Zestar Media Intelligence", "Risikan Media Zestar"), 930, 690, 286, 20, size=12, bold=True, color=C["muted"], align="right")


def _metric(slide: Any, left: float, label: str, value: str, detail: str) -> None:
    _text(slide, label.upper(), left, 170, 250, 24, size=16, bold=True, color=C["muted"])
    _text(slide, value, left, 204, 250, 58, size=42, bold=True, color=C["orange_dark"])
    _text(slide, detail, left, 270, 250, 42, size=16, color=C["muted"])


def _style_chart(chart: Any, color: str, has_legend: bool = False) -> None:
    chart.has_legend = has_legend
    if has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    chart.chart_style = 10
    chart.has_title = False
    chart.value_axis.minimum_scale = 0
    chart.value_axis.tick_labels.font.name = "Aptos"
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.name = "Aptos"
    chart.category_axis.tick_labels.font.size = Pt(9)
    for series in chart.series:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = _rgb(color)
        series.format.line.color.rgb = _rgb(color)


def _line_chart(slide: Any, categories: Sequence[str], values: Sequence[float], left: float, top: float, width: float, height: float, name: str, color: str) -> None:
    data = CategoryChartData()
    data.categories = list(categories) or [""]
    data.add_series(name, list(values) or [0])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(left * PX), Inches(top * PX), Inches(width * PX), Inches(height * PX),
        data,
    ).chart
    _style_chart(chart, color)
    series = chart.series[0]
    series.format.line.width = Pt(2.5)
    series.marker.size = 6
    series.marker.format.fill.solid()
    series.marker.format.fill.fore_color.rgb = _rgb(color)
    series.marker.format.line.color.rgb = _rgb(color)


def _bar_chart(slide: Any, categories: Sequence[str], values: Sequence[float], left: float, top: float, width: float, height: float, name: str, color: str, horizontal: bool = False, percent: bool = False) -> None:
    data = CategoryChartData()
    data.categories = list(categories) or [""]
    data.add_series(name, list(values) or [0])
    chart_type = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    chart = slide.shapes.add_chart(
        chart_type,
        Inches(left * PX), Inches(top * PX), Inches(width * PX), Inches(height * PX),
        data,
    ).chart
    _style_chart(chart, color)
    chart.plots[0].gap_width = 48
    chart.plots[0].has_data_labels = True
    labels = chart.plots[0].data_labels
    labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    labels.font.name = "Aptos"
    labels.font.size = Pt(9)
    if percent:
        chart.value_axis.tick_labels.number_format = '0.0"%"'
        labels.number_format = '0.0"%"'


def _top(items: Iterable[Dict[str, Any]], limit: int) -> List[Dict[str, Any]]:
    return list(items or [])[:limit]


def _normalize_chart_axis_ids(content: bytes) -> bytes:
    """Keep chart axis identifiers within the unsigned range required by OOXML."""
    source = BytesIO(content)
    destination = BytesIO()
    with ZipFile(source, "r") as archive, ZipFile(destination, "w", ZIP_DEFLATED) as output:
        for item in archive.infolist():
            payload = archive.read(item.filename)
            if item.filename.startswith("ppt/charts/chart") and item.filename.endswith(".xml"):
                payload = re.sub(
                    rb'(<c:(?:axId|crossAx)\s+val=")-(\d+)(")',
                    rb"\1\2\3",
                    payload,
                )
            output.writestr(item, payload)
    return destination.getvalue()


def build_pptx_report(dataset: Dict[str, Any], options: Dict[str, Any]) -> bytes:
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH)
    presentation.slide_height = Inches(SLIDE_HEIGHT)
    blank = presentation.slide_layouts[6]

    malay = str(options.get("language") or "en").lower() == "ms"
    t = lambda english, bahasa: bahasa if malay else english
    sections = set(options.get("sections") or [])
    metrics = dataset["metrics"]
    totals = metrics["totals"]
    total_engagement = sum(_safe_number(value) for value in totals.get("engagement", {}).values())
    sentiment_total = sum(_safe_number(value) for value in totals.get("sentiment", {}).values()) or 1
    sentiment_shares = {
        key: _safe_number(value) / sentiment_total * 100
        for key, value in totals.get("sentiment", {}).items()
    }
    for key in ("positive", "neutral", "negative"):
        sentiment_shares.setdefault(key, 0.0)
    dominant = max(sentiment_shares, key=sentiment_shares.get)
    days = metrics.get("days", [])
    peak = max(days, key=lambda item: sum(_safe_number(value) for value in item.get("engagement", {}).values()), default=None)
    top_source = metrics.get("sources", [None])[0] if metrics.get("sources") else None

    slide = presentation.slides.add_slide(blank)
    _set_background(slide, C["surface"])
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 790, 720, C["orange_dark"])
    _shape(slide, MSO_SHAPE.RECTANGLE, 790, 0, 490, 720, C["gold"])
    _text(slide, t("MEDIA INTELLIGENCE", "RISIKAN MEDIA"), 68, 76, 500, 30, size=18, bold=True, color="FFD9BD")
    _text(slide, options.get("reportTitle") or t("Media Intelligence Report", "Laporan Risikan Media"), 68, 160, 650, 180, size=54, bold=True, color=C["white"], valign="middle")
    _text(slide, dataset["project"]["name"], 68, 374, 600, 46, size=28, bold=True, color=C["white"])
    _text(slide, f"{dataset['dateRange']['from']}  –  {dataset['dateRange']['to']}", 68, 438, 460, 32, size=20, color="FFE8D7")
    _text(slide, options.get("organization") or t("Prepared for client review", "Disediakan untuk semakan klien"), 850, 246, 330, 100, size=28, bold=True, align="center", valign="middle")
    _text(slide, t("Live monitoring snapshot", "Petikan pemantauan langsung"), 850, 370, 330, 34, size=18, align="center")
    _text(slide, t("Zestar Media Intelligence", "Risikan Media Zestar"), 850, 628, 330, 30, size=16, bold=True, align="center")

    page = 1
    if "overview" in sections:
        page += 1
        slide = presentation.slides.add_slide(blank)
        _header(slide, t("The conversation at a glance", "Perbualan secara ringkas"), t("Executive overview", "Ringkasan eksekutif"), page)
        _metric(slide, 70, t("Mentions", "Sebutan"), _compact(totals.get("mentions")), t("Across the reporting period", "Sepanjang tempoh laporan"))
        _metric(slide, 360, t("Estimated reach", "Anggaran capaian"), _compact(totals.get("reach")), t("Potential audience exposure", "Potensi pendedahan khalayak"))
        _metric(slide, 650, t("Engagement", "Penglibatan"), _compact(total_engagement), t("Likes, comments and shares", "Suka, komen dan perkongsian"))
        _metric(slide, 940, t("Peak day", "Hari puncak"), _short_date(peak["date"], malay) if peak else "—", t("Highest daily engagement", "Penglibatan harian tertinggi"))
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 70, 360, 1140, 220, C["soft"], C["line"])
        _text(slide, t("What stands out", "Perkara utama"), 100, 390, 300, 34, size=24, bold=True, color=C["orange_dark"])
        if top_source:
            dominant_label = {"positive": t("Positive", "Positif"), "negative": t("Negative", "Negatif"), "neutral": "Neutral"}[dominant]
            message = t(
                f"{dominant_label} sentiment accounts for {sentiment_shares[dominant]:.1f}% of classified mentions. {top_source['source']} leads source volume with {_compact(top_source['mentions'])} mentions.",
                f"Sentimen {dominant_label.lower()} merangkumi {sentiment_shares[dominant]:.1f}% sebutan terkelas. {top_source['source']} mendahului jumlah sumber dengan {_compact(top_source['mentions'])} sebutan.",
            )
        else:
            message = t("Conversation data is available, but no source breakdown was returned.", "Data perbualan tersedia, tetapi tiada pecahan sumber diterima.")
        _text(slide, message, 100, 444, 1030, 92, size=24, valign="middle")
        _footer(slide, dataset, page, t)

    if "daily" in sections:
        page += 1
        slide = presentation.slides.add_slide(blank)
        _header(slide, t("Attention moved through distinct daily peaks", "Puncak harian menunjukkan perubahan perhatian"), t("Daily timeline", "Garis masa harian"), page)
        _text(slide, t("Mentions", "Sebutan"), 70, 150, 300, 34, size=24, bold=True, color=C["orange_dark"])
        _text(slide, t("Engagement", "Penglibatan"), 660, 150, 300, 34, size=24, bold=True, color=C["gold"])
        categories = [_short_date(item.get("date", ""), malay) for item in days]
        _line_chart(slide, categories, [_safe_number(item.get("mentions")) for item in days], 70, 195, 530, 365, t("Mentions", "Sebutan"), C["orange_dark"])
        _line_chart(slide, categories, [sum(_safe_number(value) for value in item.get("engagement", {}).values()) for item in days], 660, 195, 550, 365, t("Engagement", "Penglibatan"), C["gold"])
        _footer(slide, dataset, page, t)

    if "sources" in sections:
        page += 1
        slide = presentation.slides.add_slide(blank)
        _header(slide, t("A small group of sources drove most visibility", "Beberapa sumber memacu keterlihatan"), t("Source performance", "Prestasi sumber"), page)
        source_items = _top(metrics.get("sources", []), 8)
        _bar_chart(slide, [_source_label(item.get("source"), malay) for item in source_items], [_safe_number(item.get("mentions")) for item in source_items], 70, 158, 760, 450, t("Mentions", "Sebutan"), C["orange"])
        _text(slide, t("Leading source", "Sumber utama"), 890, 200, 280, 28, size=18, bold=True, color=C["muted"])
        _text(slide, top_source.get("source") if top_source else "—", 890, 245, 280, 70, size=34, bold=True, color=C["orange_dark"], valign="middle")
        detail = f"{_compact(top_source.get('mentions'))} {t('mentions', 'sebutan')}\n{_compact(top_source.get('reach'))} {t('estimated reach', 'anggaran capaian')}" if top_source else "—"
        _text(slide, detail, 890, 335, 280, 90, size=20)
        _text(slide, t("Source reach is aggregated; it is not a unique-audience count.", "Capaian sumber ialah agregat dan bukan kiraan khalayak unik."), 890, 470, 280, 80, size=16, color=C["muted"])
        _footer(slide, dataset, page, t)

    if "overview" in sections:
        page += 1
        slide = presentation.slides.add_slide(blank)
        dominant_label = {"positive": t("Positive", "Positif"), "negative": t("Negative", "Negatif"), "neutral": "Neutral"}[dominant]
        _header(slide, t(f"{dominant_label} sentiment formed the largest share", f"Sentimen {dominant_label.lower()} mendominasi perbualan"), t("Sentiment", "Sentimen"), page)
        dominant_color = C[dominant]
        _text(slide, f"{sentiment_shares[dominant]:.1f}%", 90, 190, 540, 100, size=72, bold=True, color=dominant_color, align="center")
        _text(slide, t(f"{dominant_label} share of classified mentions", f"Bahagian sebutan terkelas: {dominant_label.lower()}"), 110, 295, 500, 42, size=23, bold=True, align="center")
        segment_left = 105.0
        for key in ("positive", "neutral", "negative"):
            segment_width = 500 * sentiment_shares[key] / 100
            if segment_width > 0:
                _shape(slide, MSO_SHAPE.RECTANGLE, segment_left, 375, segment_width, 54, C[key])
            segment_left += segment_width
        _text(slide, t("Share of classified mentions", "Bahagian sebutan terkelas"), 105, 452, 500, 28, size=16, color=C["muted"], align="center")
        sentiment_rows = [(t("Positive", "Positif"), "positive"), ("Neutral", "neutral"), (t("Negative", "Negatif"), "negative")]
        for index, (label, key) in enumerate(sentiment_rows):
            top = 195 + index * 120
            _shape(slide, MSO_SHAPE.RECTANGLE, 760, top, 8, 76, C[key])
            _text(slide, label, 790, top, 220, 30, size=22, bold=True)
            _text(slide, f"{sentiment_shares[key]:.1f}%", 1030, top - 5, 150, 38, size=30, bold=True, color=C[key], align="right")
            _text(slide, f"{_compact(totals.get('sentiment', {}).get(key))} {t('mentions', 'sebutan')}", 790, top + 41, 390, 26, size=16, color=C["muted"])
        _footer(slide, dataset, page, t)

    if "topics" in sections:
        page += 1
        slide = presentation.slides.add_slide(blank)
        topics = _top(dataset.get("topics", {}).get("items", []), 6)
        _header(slide, t("The leading themes reveal where attention clustered", "Tema utama menunjukkan tumpuan perhatian"), t("Topics and share of voice", "Topik dan bahagian suara"), page)
        _bar_chart(slide, [_truncate(item.get("name"), 32) for item in topics], [_safe_number(item.get("shareOfVoice")) for item in topics], 70, 160, 660, 430, t("Share of voice", "Bahagian suara"), C["orange"], horizontal=True, percent=True)
        for index, topic in enumerate(topics[:3]):
            top = 175 + index * 138
            _text(slide, str(index + 1), 785, top, 40, 40, size=20, bold=True, color=C["orange_dark"], align="center")
            _text(slide, topic.get("name", ""), 845, top, 360, 34, size=21, bold=True)
            _text(slide, _truncate(topic.get("description"), 150), 845, top + 42, 360, 74, size=16, color=C["muted"])
        _footer(slide, dataset, page, t)

    if sections.intersection({"sources", "authors"}):
        page += 1
        slide = presentation.slides.add_slide(blank)
        _header(slide, t("Visibility came from both publishers and influential voices", "Penerbit dan suara utama memacu keterlihatan"), t("Who and where", "Siapa dan di mana"), page)
        _text(slide, t("Leading domains", "Domain utama"), 70, 155, 500, 40, size=26, bold=True, color=C["orange_dark"])
        _text(slide, t("Influential authors", "Pengarang berpengaruh"), 680, 155, 500, 40, size=26, bold=True, color=C["orange_dark"])
        domains = _top(dataset.get("sources", {}).get("domains", {}).get("items", []), 6)
        authors = _top(dataset.get("sources", {}).get("authors", {}).get("items", []), 6)
        for index, item in enumerate(domains):
            top = 220 + index * 65
            _text(slide, str(index + 1), 70, top, 36, 28, size=18, bold=True, color=C["orange_dark"])
            _text(slide, item.get("domain", ""), 120, top, 300, 28, size=18, bold=True)
            _text(slide, f"{_compact(item.get('mentions'))} {t('mentions', 'sebutan')}  •  {_compact(item.get('reach'))} {t('reach', 'capaian')}", 120, top + 30, 440, 24, size=14, color=C["muted"])
        for index, item in enumerate(authors):
            top = 220 + index * 65
            _text(slide, str(index + 1), 680, top, 36, 28, size=18, bold=True, color=C["gold"])
            _text(slide, _truncate(item.get("name"), 34), 730, top, 300, 28, size=18, bold=True)
            _text(slide, f"{_compact(item.get('followers'))} {t('followers', 'pengikut')}  •  {_compact(item.get('reach'))} {t('reach', 'capaian')}", 730, top + 30, 440, 24, size=14, color=C["muted"])
        _footer(slide, dataset, page, t)

    if "mentions" in sections:
        page += 1
        slide = presentation.slides.add_slide(blank)
        _header(slide, t("Representative mentions capture how coverage was framed", "Sebutan wakil menunjukkan bingkai liputan"), t("Featured coverage", "Liputan pilihan"), page)
        candidates = [item for item in dataset.get("mentions", {}).get("items", []) if item.get("title") or item.get("content")][:3]
        for index, item in enumerate(candidates):
            left = 64 + index * 400
            _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, 165, 360, 430, C["soft"] if index == 0 else C["white"], C["line"])
            sentiment = str(item.get("sentiment") or "neutral").lower()
            sentiment_color = C.get(sentiment, C["neutral"])
            _text(slide, f"{item.get('category', '')}  •  {sentiment}", left + 24, 192, 312, 28, size=16, bold=True, color=sentiment_color)
            _text(slide, _truncate(item.get("title") or t("Untitled mention", "Sebutan tanpa tajuk"), 90), left + 24, 240, 312, 92, size=22, bold=True)
            excerpt = item.get("content") or item.get("restrictionReason") or t("No excerpt supplied by the data source.", "Tiada petikan dibekalkan oleh sumber data.")
            _text(slide, _truncate(excerpt, 250), left + 24, 350, 312, 150, size=16, color=C["muted"])
            _text(slide, f"{item.get('host') or t('Source unavailable', 'Sumber tidak tersedia')}\n{item.get('date', '')}", left + 24, 525, 312, 48, size=14, bold=True, color=C["orange_dark"])
        if not candidates:
            _text(slide, t("No unrestricted mention excerpts were available for this period.", "Tiada petikan sebutan tanpa sekatan tersedia bagi tempoh ini."), 170, 300, 940, 70, size=26, color=C["muted"], align="center")
        _footer(slide, dataset, page, t)

    if sections.intersection({"links", "hashtags"}):
        page += 1
        slide = presentation.slides.add_slide(blank)
        _header(slide, t("Shared links and hashtags reveal the conversation signals", "Pautan dan tanda pagar menunjukkan isyarat perbualan"), t("Conversation signals", "Isyarat perbualan"), page)
        _text(slide, t("Trending links", "Pautan sohor kini"), 70, 155, 520, 40, size=26, bold=True, color=C["orange_dark"])
        _text(slide, t("Trending hashtags", "Tanda pagar sohor kini"), 700, 155, 490, 40, size=26, bold=True, color=C["orange_dark"])
        for index, item in enumerate(_top(dataset.get("sources", {}).get("links", {}).get("items", []), 6)):
            top = 220 + index * 62
            _text(slide, _truncate(item.get("url"), 62), 70, top, 470, 30, size=17, bold=True)
            _text(slide, f"{_compact(item.get('mentions'))} {t('mentions', 'sebutan')}", 550, top, 100, 26, size=16, bold=True, color=C["orange_dark"], align="right")
        for index, item in enumerate(_top(dataset.get("sources", {}).get("hashtags", {}).get("items", []), 8)):
            top = 220 + index * 48
            _text(slide, item.get("hashtag", ""), 700, top, 310, 28, size=19, bold=True, color=C["orange_dark"])
            _text(slide, f"{_compact(item.get('mentions'))}  •  {_compact(item.get('reach'))} {t('reach', 'capaian')}", 1020, top, 170, 26, size=15, color=C["muted"], align="right")
        _footer(slide, dataset, page, t)

    page += 1
    slide = presentation.slides.add_slide(blank)
    _header(slide, t("Read the findings within the limits of the available data", "Tafsir dapatan berdasarkan had data yang tersedia"), t("Methodology", "Metodologi"), page)
    notes = [
        t("Metrics are generated from a live monitoring snapshot for the selected project and reporting period.", "Metrik dijana daripada petikan pemantauan langsung bagi projek dan tempoh laporan yang dipilih."),
        t("Reach represents potential exposure and should not be interpreted as a unique-audience count.", "Capaian mewakili potensi pendedahan dan bukan kiraan khalayak unik."),
        t("Individual mentions do not include per-post reach, likes, comments or shares; featured coverage is representative, not performance-ranked.", "Sebutan individu tidak mempunyai capaian, suka, komen atau perkongsian setiap siaran; liputan pilihan ialah wakil dan bukan kedudukan prestasi."),
        t("Facebook, Instagram and X may restrict post text or source links in API responses.", "Facebook, Instagram dan X mungkin mengehadkan teks siaran atau pautan sumber dalam respons API."),
        t(f"Raw mentions were capped at {int(dataset.get('mentions', {}).get('limit', 0)):,} rows.", f"Sebutan mentah dihadkan kepada {int(dataset.get('mentions', {}).get('limit', 0)):,} baris.") if dataset.get("mentions", {}).get("truncated") else t("All mention pages returned by the monitoring service for this period were included.", "Semua halaman sebutan yang dikembalikan oleh perkhidmatan pemantauan bagi tempoh ini telah disertakan."),
    ]
    for index, note in enumerate(notes):
        top = 170 + index * 92
        _text(slide, str(index + 1), 90, top, 48, 40, size=24, bold=True, color=C["orange_dark"], align="center")
        _text(slide, note, 165, top, 980, 70, size=20, valign="middle")
    _footer(slide, dataset, page, t)

    output = BytesIO()
    presentation.save(output)
    return _normalize_chart_axis_ids(output.getvalue())
